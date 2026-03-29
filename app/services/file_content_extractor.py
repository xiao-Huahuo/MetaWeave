from io import BytesIO
from pathlib import Path
from typing import Optional, Iterable, Iterator, Union

from docx import Document
from docx.oxml.ns import qn
from docx.text.paragraph import Paragraph
from docx.table import Table
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz

from app.core.logger import global_logger as logger
from app.services.ocr_service import OCRService


class FileContentExtractor:
    TEXT_EXTENSIONS = {".txt", ".md"}
    SPREADSHEET_EXTENSIONS = {".xlsx", ".xls"}
    DOC_EXTENSIONS = {".doc", ".docx"}
    PDF_EXTENSIONS = {".pdf"}
    PPT_EXTENSIONS = {".ppt", ".pptx"}
    IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp"}

    @classmethod
    def is_supported(cls, file_path: str) -> bool:
        suffix = Path(file_path).suffix.lower()
        return suffix in (
            cls.TEXT_EXTENSIONS
            | cls.SPREADSHEET_EXTENSIONS
            | cls.DOC_EXTENSIONS
            | cls.PDF_EXTENSIONS
            | cls.PPT_EXTENSIONS
            | cls.IMAGE_EXTENSIONS
        )

    @classmethod
    def extract(cls, file_path: str) -> Optional[str]:
        """
        统一入口：解析任意格式文件内容。
        """
        try:
            path = Path(file_path)
            suffix = path.suffix.lower()

            if suffix in cls.TEXT_EXTENSIONS:
                return cls._read_text_file(file_path)
            if suffix in cls.SPREADSHEET_EXTENSIONS:
                return cls._read_spreadsheet_file(file_path, suffix)
            if suffix in cls.DOC_EXTENSIONS:
                return cls._read_docx_file(file_path)
            if suffix in cls.PDF_EXTENSIONS:
                return cls._read_pdf_file(file_path)
            if suffix in cls.PPT_EXTENSIONS:
                return cls._read_pptx_file(file_path)
            if suffix in cls.IMAGE_EXTENSIONS:
                return OCRService.recognize_image_path(file_path)

            return None
        except Exception as e:
            logger.error(f"提取文件内容失败 {file_path}: {e}")
            return None

    @staticmethod
    def _read_text_file(file_path: str) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gbk", "gb2312", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding, errors="strict") as f:
                    return f.read()
            except Exception:
                continue
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    def _read_spreadsheet_file(file_path: str, suffix: str) -> Optional[str]:
        try:
            if suffix == ".xlsx":
                return FileContentExtractor._read_xlsx(file_path)
            if suffix == ".xls":
                return FileContentExtractor._read_xls(file_path)
        except Exception as e:
            logger.error(f"读取表格失败 {file_path}: {e}")
        return None

    @staticmethod
    def _read_xlsx(file_path: str) -> Optional[str]:
        try:
            import openpyxl
        except Exception as e:
            logger.error(f"缺少 openpyxl, 无法读取 xlsx: {e}")
            return None

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"[Sheet] {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = FileContentExtractor._join_row_values(row)
                if row_text:
                    lines.append(row_text)
        wb.close()
        return "\n".join(lines) if lines else None

    @staticmethod
    def _read_xls(file_path: str) -> Optional[str]:
        try:
            import xlrd
        except Exception as e:
            logger.error(f"缺少 xlrd, 无法读取 xls: {e}")
            return None

        book = xlrd.open_workbook(file_path)
        lines = []
        for sheet in book.sheets():
            lines.append(f"[Sheet] {sheet.name}")
            for r in range(sheet.nrows):
                row = sheet.row_values(r)
                row_text = FileContentExtractor._join_row_values(row)
                if row_text:
                    lines.append(row_text)
        return "\n".join(lines) if lines else None

    @staticmethod
    def _join_row_values(values: Iterable) -> str:
        parts = []
        for v in values:
            if v is None:
                parts.append("")
            else:
                parts.append(str(v))
        row_text = "\t".join(parts).strip()
        return row_text

    @classmethod
    def _read_docx_file(cls, file_path: str) -> Optional[str]:
        if Path(file_path).suffix.lower() == ".doc":
            logger.error(f"暂不支持解析 .doc 文件: {file_path}")
            return None
        doc = Document(file_path)
        parts = []
        for block in cls._iter_block_items(doc):
            if isinstance(block, Paragraph):
                parts.extend(cls._extract_paragraph_runs(doc, block))
            elif isinstance(block, Table):
                for row in block.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            parts.extend(cls._extract_paragraph_runs(doc, paragraph))
        text = "\n".join([p for p in parts if p])
        return text if text else None

    @staticmethod
    def _iter_block_items(parent) -> Iterator[Union[Paragraph, Table]]:
        for child in parent.element.body.iterchildren():
            if child.tag.endswith("}p"):
                yield Paragraph(child, parent)
            elif child.tag.endswith("}tbl"):
                yield Table(child, parent)

    @staticmethod
    def _extract_paragraph_runs(doc: Document, paragraph: Paragraph) -> Iterable[str]:
        parts = []
        for run in paragraph.runs:
            if run.text:
                parts.append(run.text)
            blips = run._element.xpath(".//a:blip")
            for blip in blips:
                r_id = blip.get(qn("r:embed"))
                image_part = doc.part.related_parts.get(r_id)
                if image_part and image_part.blob:
                    ocr_text = OCRService.recognize_image_bytes(image_part.blob)
                    if ocr_text:
                        parts.append(ocr_text)
        return parts

    @classmethod
    def _read_pdf_file(cls, file_path: str) -> Optional[str]:
        doc = fitz.open(file_path)
        parts = []
        for page in doc:
            text = page.get_text().strip()
            if len(text) >= 10:
                parts.append(text)
                continue
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            img_bytes = pix.tobytes("png")
            ocr_text = OCRService.recognize_image_bytes(img_bytes)
            if ocr_text:
                parts.append(ocr_text)
        doc.close()
        text = "\n".join([p for p in parts if p])
        return text if text else None

    @classmethod
    def _read_pptx_file(cls, file_path: str) -> Optional[str]:
        if Path(file_path).suffix.lower() == ".ppt":
            logger.error(f"暂不支持解析 .ppt 文件: {file_path}")
            return None
        prs = Presentation(file_path)
        parts = []
        for index, slide in enumerate(prs.slides, start=1):
            parts.append(f"[Slide {index}]")
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text:
                        parts.append(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    if image and image.blob:
                        ocr_text = OCRService.recognize_image_bytes(image.blob)
                        if ocr_text:
                            parts.append(ocr_text)
        text = "\n".join([p for p in parts if p])
        return text if text else None
